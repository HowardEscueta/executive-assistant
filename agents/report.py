"""
Report Generator
Creates a PowerPoint briefing deck with New Songs, AI Automation Trends, and Freelance Leads.
Dark theme, professional styling -- same vibe as the youtube-trends project.
"""

import os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
REPORTS_DIR = os.path.join(PROJECT_ROOT, "output", "reports")
os.makedirs(REPORTS_DIR, exist_ok=True)

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xD2, 0x7A)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_hyperlink(slide, text, url, left, top, width, height, font_size=9, color=ACCENT_BLUE):
    """Add a clickable hyperlink to a slide."""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.underline = True
    # Add hyperlink via oxml
    r_elem = run._r
    rPr = r_elem.get_or_add_rPr()
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {})
    rPr.append(hlinkClick)
    # Create relationship and set rId
    rel_type = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'
    slide_part = slide.part
    rId = slide_part.relate_to(url, rel_type, is_external=True)
    hlinkClick.set(qn('r:id'), rId)
    return tf


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)
    add_text(slide, "Morning Briefing", 0.5, 1.5, 9, 1, font_size=36, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, date.today().isoformat(), 0.5, 2.5, 9, 0.5, font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, "New Songs  |  AI Automation  |  Freelance Leads", 0.5, 3.2, 9, 0.5, font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def add_section_header(prs, title, color=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, title, 0.5, 2.5, 9, 1, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    return slide


def add_songs_slides(prs, songs):
    if not songs:
        return
    add_section_header(prs, "New Songs (No Lyrics)", ACCENT_RED)

    # 5 songs per slide
    for chunk_start in range(0, len(songs), 5):
        chunk = songs[chunk_start:chunk_start + 5]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_text(slide, "New Songs", 0.3, 0.2, 5, 0.4, font_size=18, color=ACCENT_RED, bold=True)

        y = 0.7
        for i, s in enumerate(chunk, chunk_start + 1):
            # Title
            add_text(slide, f"{i}. {s['title']}", 0.3, y, 9.2, 0.3, font_size=13, color=WHITE, bold=True)
            y += 0.3
            # Details
            details = f"{s['channel']}  |  {s['views']:,} views  |  {s['upload_date']}  |  {s['duration_min']}"
            add_text(slide, details, 0.5, y, 9, 0.25, font_size=10, color=LIGHT_GRAY)
            y += 0.25
            # Clickable URL
            add_hyperlink(slide, s['url'], s['url'], 0.5, y, 9, 0.25)
            y += 0.35


def add_ai_slides(prs, videos):
    if not videos:
        return
    add_section_header(prs, "AI Automation Trends", ACCENT_GREEN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, "AI Automation Trends", 0.3, 0.2, 5, 0.4, font_size=18, color=ACCENT_GREEN, bold=True)

    y = 0.7
    for i, v in enumerate(videos[:10], 1):
        add_text(slide, f"{i}. {v['title']}", 0.3, y, 9.2, 0.3, font_size=13, color=WHITE, bold=True)
        y += 0.3
        details = f"{v['channel']}  |  {v['views']:,} views  |  Engagement: {v['engagement_rate']}%  |  {v['upload_date']}"
        add_text(slide, details, 0.5, y, 9, 0.25, font_size=10, color=LIGHT_GRAY)
        y += 0.25
        add_hyperlink(slide, v['url'], v['url'], 0.5, y, 9, 0.25)
        y += 0.35


def add_leads_slides(prs, leads):
    if not leads:
        return
    hot = [l for l in leads if l.get('score', 0) >= 3]
    others = [l for l in leads if l.get('score', 0) < 3]

    add_section_header(prs, f"Freelance Leads ({len(hot)} Hot)", ACCENT_YELLOW)

    if hot:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_text(slide, "Hot Leads", 0.3, 0.2, 5, 0.4, font_size=18, color=ACCENT_YELLOW, bold=True)

        y = 0.7
        for i, l in enumerate(hot, 1):
            source = l.get('source', 'Web')
            budget = l.get('budget', '')
            title_line = f"{i}. [{source}] {l['title']}"
            if budget:
                title_line += f"  ({budget})"
            add_text(slide, title_line, 0.3, y, 9.2, 0.3, font_size=13, color=WHITE, bold=True)
            y += 0.3
            if l.get('snippet'):
                add_text(slide, l['snippet'][:120], 0.5, y, 9, 0.25, font_size=10, color=LIGHT_GRAY)
                y += 0.25
            add_hyperlink(slide, l['url'], l['url'], 0.5, y, 9, 0.25)
            y += 0.2
            status = l.get('status', '')
            score_text = f"Score: {l.get('score', 0)}/10"
            if status:
                score_text += f"  |  Status: {status}"
            add_text(slide, score_text, 0.5, y, 4, 0.2, font_size=10, color=ACCENT_YELLOW, bold=True)
            y += 0.35

    if others:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_text(slide, "Other Leads", 0.3, 0.2, 5, 0.4, font_size=18, color=MED_GRAY, bold=True)

        y = 0.7
        for i, l in enumerate(others, 1):
            add_text(slide, f"{i}. {l['title'][:70]}", 0.3, y, 9.2, 0.25, font_size=11, color=LIGHT_GRAY)
            y += 0.25
            add_hyperlink(slide, l['url'], l['url'], 0.5, y, 9, 0.25)
            y += 0.3


def add_local_leads_slides(prs, local_leads):
    if not local_leads:
        return
    ACCENT_ORANGE = RGBColor(0xFF, 0x7A, 0x00)
    add_section_header(prs, f"Local Business Leads ({len(local_leads)} No Website)", ACCENT_ORANGE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, "Local Businesses Without Websites", 0.3, 0.2, 9, 0.4, font_size=18, color=ACCENT_ORANGE, bold=True)

    y = 0.7
    for i, b in enumerate(local_leads[:8], 1):
        if y > 4.8:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide)
            add_text(slide, "Local Businesses (cont.)", 0.3, 0.2, 9, 0.4, font_size=18, color=ACCENT_ORANGE, bold=True)
            y = 0.7

        add_text(slide, f"{i}. {b['name']}  [{b['category'].title()}]", 0.3, y, 9.2, 0.3, font_size=13, color=WHITE, bold=True)
        y += 0.28
        rating_str = f"{b['rating']}/5 ({b['user_ratings_total']} reviews)" if b.get("rating") else "No ratings"
        details = f"{b['address']}  |  {b.get('phone') or 'No phone'}  |  {rating_str}"
        add_text(slide, details, 0.5, y, 9, 0.22, font_size=9, color=LIGHT_GRAY)
        y += 0.22
        if b.get("maps_url"):
            add_hyperlink(slide, "View on Google Maps", b["maps_url"], 0.5, y, 4, 0.22)
        y += 0.35


def generate_report(new_songs=None, ai_trends=None, leads=None, local_leads=None):
    """Generate a PowerPoint briefing report. Returns the file path."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    add_title_slide(prs)
    add_songs_slides(prs, new_songs)
    add_ai_slides(prs, ai_trends)
    add_leads_slides(prs, leads)
    add_local_leads_slides(prs, local_leads)

    today = date.today().isoformat()
    filepath = os.path.join(REPORTS_DIR, f"{today}_morning_briefing.pptx")
    prs.save(filepath)
    print(f"Report saved: {filepath}")
    return filepath


if __name__ == "__main__":
    print("Run via morning_briefing.py to generate a report with live data.")
